"""
Document Processing Service for MEMSHADOW

Handles extraction, parsing, and processing of multiple document formats:
- PDF (text, images, metadata)
- Word documents (DOCX)
- PowerPoint (PPTX)
- Excel (XLSX)
- HTML/Markdown
- Plain text
- OCR for scanned documents
"""

import io
import hashlib
import magic
from typing import Dict, Any, List, Optional, Tuple, BinaryIO
from pathlib import Path
from datetime import datetime
import structlog

import fitz  # PyMuPDF
from PIL import Image
import pytesseract
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
from bs4 import BeautifulSoup
import markdown

from app.core.config import settings

logger = structlog.get_logger()


class DocumentProcessingError(Exception):
    """Custom exception for document processing errors"""
    pass


class DocumentProcessor:
    """
    Main document processing class that handles multiple file formats
    """

    # Supported file types and their MIME types
    SUPPORTED_TYPES = {
        'application/pdf': 'pdf',
        'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
        'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
        'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
        'text/html': 'html',
        'text/markdown': 'markdown',
        'text/plain': 'txt',
        'image/png': 'image',
        'image/jpeg': 'image',
        'image/jpg': 'image',
        'image/tiff': 'image',
    }

    # Maximum file size: 50MB
    MAX_FILE_SIZE = 50 * 1024 * 1024

    # Maximum pages for PDFs to prevent abuse
    MAX_PDF_PAGES = 500

    def __init__(self):
        self.magic_mime = magic.Magic(mime=True)

    async def process_document(
        self,
        file_content: bytes,
        filename: str,
        user_id: str
    ) -> Dict[str, Any]:
        """
        Main entry point for document processing

        Args:
            file_content: Binary content of the file
            filename: Original filename
            user_id: ID of user uploading the document

        Returns:
            Dictionary containing extracted text, images, metadata
        """
        try:
            # Validate file size
            if len(file_content) > self.MAX_FILE_SIZE:
                raise DocumentProcessingError(
                    f"File size exceeds maximum allowed size of {self.MAX_FILE_SIZE} bytes"
                )

            # Detect file type
            mime_type = self.magic_mime.from_buffer(file_content)
            file_type = self.SUPPORTED_TYPES.get(mime_type)

            if not file_type:
                raise DocumentProcessingError(
                    f"Unsupported file type: {mime_type}"
                )

            logger.info(
                "Processing document",
                filename=filename,
                mime_type=mime_type,
                file_type=file_type,
                size=len(file_content),
                user_id=user_id
            )

            # Route to appropriate processor
            if file_type == 'pdf':
                result = await self._process_pdf(file_content)
            elif file_type == 'docx':
                result = await self._process_docx(file_content)
            elif file_type == 'pptx':
                result = await self._process_pptx(file_content)
            elif file_type == 'xlsx':
                result = await self._process_xlsx(file_content)
            elif file_type == 'html':
                result = await self._process_html(file_content)
            elif file_type == 'markdown':
                result = await self._process_markdown(file_content)
            elif file_type == 'txt':
                result = await self._process_text(file_content)
            elif file_type == 'image':
                result = await self._process_image(file_content)
            else:
                raise DocumentProcessingError(f"No processor for type: {file_type}")

            # Add common metadata
            result['document_metadata'] = {
                'original_filename': filename,
                'mime_type': mime_type,
                'file_type': file_type,
                'file_size': len(file_content),
                'processed_at': datetime.utcnow().isoformat(),
                'content_hash': hashlib.sha256(file_content).hexdigest(),
                'user_id': user_id
            }

            logger.info(
                "Document processed successfully",
                filename=filename,
                text_length=len(result.get('text', '')),
                num_chunks=len(result.get('chunks', [])),
                num_images=len(result.get('images', []))
            )

            return result

        except Exception as e:
            logger.error(
                "Document processing failed",
                filename=filename,
                error=str(e),
                error_type=type(e).__name__
            )
            raise DocumentProcessingError(f"Failed to process document: {str(e)}")

    async def _process_pdf(self, content: bytes) -> Dict[str, Any]:
        """
        Process PDF documents with advanced text and image extraction
        """
        try:
            doc = fitz.open(stream=content, filetype="pdf")

            # Check page count
            if doc.page_count > self.MAX_PDF_PAGES:
                doc.close()
                raise DocumentProcessingError(
                    f"PDF has too many pages ({doc.page_count}). Maximum allowed: {self.MAX_PDF_PAGES}"
                )

            text_content = []
            images = []
            chunks = []

            # Extract text and images from each page
            for page_num in range(doc.page_count):
                page = doc[page_num]

                # Extract text
                page_text = page.get_text("text")

                # If page has minimal text, try OCR on page image
                if len(page_text.strip()) < 50:
                    page_text = await self._ocr_pdf_page(page)

                if page_text.strip():
                    text_content.append(f"[Page {page_num + 1}]\n{page_text}")

                    # Create chunk for this page
                    chunks.append({
                        'page': page_num + 1,
                        'text': page_text,
                        'char_count': len(page_text)
                    })

                # Extract images
                image_list = page.get_images()
                for img_idx, img in enumerate(image_list):
                    try:
                        xref = img[0]
                        base_image = doc.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_ext = base_image["ext"]

                        # Optionally: Run OCR on images
                        image_text = await self._ocr_image(image_bytes)

                        images.append({
                            'page': page_num + 1,
                            'index': img_idx,
                            'format': image_ext,
                            'size': len(image_bytes),
                            'ocr_text': image_text,
                            # Store base64 or upload to storage
                            'data_hash': hashlib.md5(image_bytes).hexdigest()
                        })
                    except Exception as e:
                        logger.warning(
                            "Failed to extract image",
                            page=page_num,
                            image_idx=img_idx,
                            error=str(e)
                        )

            # Get PDF metadata
            metadata = doc.metadata

            doc.close()

            full_text = "\n\n".join(text_content)

            return {
                'text': full_text,
                'chunks': chunks,
                'images': images,
                'pdf_metadata': {
                    'title': metadata.get('title', ''),
                    'author': metadata.get('author', ''),
                    'subject': metadata.get('subject', ''),
                    'keywords': metadata.get('keywords', ''),
                    'creator': metadata.get('creator', ''),
                    'producer': metadata.get('producer', ''),
                    'creation_date': metadata.get('creationDate', ''),
                    'mod_date': metadata.get('modDate', ''),
                    'page_count': len(chunks)
                },
                'structure': {
                    'type': 'pdf',
                    'pages': len(chunks),
                    'total_images': len(images),
                    'total_chars': sum(c['char_count'] for c in chunks)
                }
            }

        except Exception as e:
            logger.error("PDF processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process PDF: {str(e)}")

    async def _ocr_pdf_page(self, page) -> str:
        """
        Perform OCR on a PDF page by converting it to an image
        """
        try:
            # Render page to image
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
            img_data = pix.tobytes("png")

            # Perform OCR
            image = Image.open(io.BytesIO(img_data))
            text = pytesseract.image_to_string(image)

            return text

        except Exception as e:
            logger.warning("OCR on PDF page failed", error=str(e))
            return ""

    async def _ocr_image(self, image_bytes: bytes) -> str:
        """
        Perform OCR on an image
        """
        try:
            image = Image.open(io.BytesIO(image_bytes))
            text = pytesseract.image_to_string(image)
            return text.strip()
        except Exception as e:
            logger.warning("Image OCR failed", error=str(e))
            return ""

    async def _process_docx(self, content: bytes) -> Dict[str, Any]:
        """
        Process Word documents (DOCX)
        """
        try:
            doc = DocxDocument(io.BytesIO(content))

            paragraphs = []
            tables = []
            images = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append({
                        'text': para.text,
                        'style': para.style.name if para.style else None
                    })

            # Extract tables
            for table_idx, table in enumerate(doc.tables):
                table_data = []
                for row in table.rows:
                    row_data = [cell.text for cell in row.cells]
                    table_data.append(row_data)
                tables.append({
                    'index': table_idx,
                    'data': table_data
                })

            # Extract images (from relationships)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_bytes = rel.target_part.blob
                        images.append({
                            'index': len(images),
                            'size': len(image_bytes),
                            'data_hash': hashlib.md5(image_bytes).hexdigest()
                        })
                    except Exception as e:
                        logger.warning("Failed to extract DOCX image", error=str(e))

            # Combine all text
            full_text = "\n\n".join(p['text'] for p in paragraphs)

            # Add table text
            for table in tables:
                table_text = "\n".join(["\t".join(row) for row in table['data']])
                full_text += f"\n\n[Table {table['index'] + 1}]\n{table_text}"

            # Get core properties
            core_props = doc.core_properties

            return {
                'text': full_text,
                'chunks': paragraphs,
                'tables': tables,
                'images': images,
                'docx_metadata': {
                    'title': core_props.title or '',
                    'author': core_props.author or '',
                    'subject': core_props.subject or '',
                    'keywords': core_props.keywords or '',
                    'created': core_props.created.isoformat() if core_props.created else '',
                    'modified': core_props.modified.isoformat() if core_props.modified else '',
                    'paragraph_count': len(paragraphs),
                    'table_count': len(tables)
                },
                'structure': {
                    'type': 'docx',
                    'paragraphs': len(paragraphs),
                    'tables': len(tables),
                    'images': len(images)
                }
            }

        except Exception as e:
            logger.error("DOCX processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process DOCX: {str(e)}")

    async def _process_pptx(self, content: bytes) -> Dict[str, Any]:
        """
        Process PowerPoint presentations (PPTX)
        """
        try:
            prs = Presentation(io.BytesIO(content))

            slides = []
            images = []

            for slide_idx, slide in enumerate(prs.slides):
                slide_text = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)

                    # Extract images
                    if shape.shape_type == 13:  # Picture
                        try:
                            image = shape.image
                            image_bytes = image.blob
                            images.append({
                                'slide': slide_idx + 1,
                                'size': len(image_bytes),
                                'content_type': image.content_type,
                                'data_hash': hashlib.md5(image_bytes).hexdigest()
                            })
                        except Exception as e:
                            logger.warning("Failed to extract PPTX image", error=str(e))

                slides.append({
                    'slide_number': slide_idx + 1,
                    'text': "\n".join(slide_text),
                    'shape_count': len(slide.shapes)
                })

            # Combine all slide text
            full_text = "\n\n".join(
                f"[Slide {s['slide_number']}]\n{s['text']}" for s in slides
            )

            return {
                'text': full_text,
                'chunks': slides,
                'images': images,
                'pptx_metadata': {
                    'slide_count': len(slides),
                    'total_shapes': sum(s['shape_count'] for s in slides),
                    'image_count': len(images)
                },
                'structure': {
                    'type': 'pptx',
                    'slides': len(slides),
                    'images': len(images)
                }
            }

        except Exception as e:
            logger.error("PPTX processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process PPTX: {str(e)}")

    async def _process_xlsx(self, content: bytes) -> Dict[str, Any]:
        """
        Process Excel spreadsheets (XLSX)
        """
        try:
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)

            sheets = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]

                rows = []
                for row in sheet.iter_rows(values_only=True):
                    # Convert row to strings and filter empty
                    row_data = [str(cell) if cell is not None else '' for cell in row]
                    if any(row_data):  # Skip completely empty rows
                        rows.append(row_data)

                # Convert to text representation
                sheet_text = "\n".join(["\t".join(row) for row in rows])

                sheets.append({
                    'name': sheet_name,
                    'rows': len(rows),
                    'text': sheet_text
                })

            # Combine all sheets
            full_text = "\n\n".join(
                f"[Sheet: {s['name']}]\n{s['text']}" for s in sheets
            )

            wb.close()

            return {
                'text': full_text,
                'chunks': sheets,
                'xlsx_metadata': {
                    'sheet_count': len(sheets),
                    'total_rows': sum(s['rows'] for s in sheets)
                },
                'structure': {
                    'type': 'xlsx',
                    'sheets': len(sheets)
                }
            }

        except Exception as e:
            logger.error("XLSX processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process XLSX: {str(e)}")

    async def _process_html(self, content: bytes) -> Dict[str, Any]:
        """
        Process HTML documents
        """
        try:
            html_text = content.decode('utf-8')
            soup = BeautifulSoup(html_text, 'lxml')

            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()

            # Extract text
            text = soup.get_text()

            # Clean up whitespace
            lines = (line.strip() for line in text.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            text = '\n'.join(chunk for chunk in chunks if chunk)

            # Extract metadata
            title = soup.title.string if soup.title else ''
            meta_description = ''
            meta_keywords = ''

            if soup.find('meta', attrs={'name': 'description'}):
                meta_description = soup.find('meta', attrs={'name': 'description'})['content']
            if soup.find('meta', attrs={'name': 'keywords'}):
                meta_keywords = soup.find('meta', attrs={'name': 'keywords'})['content']

            return {
                'text': text,
                'html_metadata': {
                    'title': title,
                    'description': meta_description,
                    'keywords': meta_keywords
                },
                'structure': {
                    'type': 'html'
                }
            }

        except Exception as e:
            logger.error("HTML processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process HTML: {str(e)}")

    async def _process_markdown(self, content: bytes) -> Dict[str, Any]:
        """
        Process Markdown documents
        """
        try:
            md_text = content.decode('utf-8')

            # Convert to HTML first for structure
            html = markdown.markdown(md_text)
            soup = BeautifulSoup(html, 'html.parser')

            # Extract plain text
            text = soup.get_text()

            return {
                'text': text,
                'raw_markdown': md_text,
                'structure': {
                    'type': 'markdown'
                }
            }

        except Exception as e:
            logger.error("Markdown processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process Markdown: {str(e)}")

    async def _process_text(self, content: bytes) -> Dict[str, Any]:
        """
        Process plain text documents
        """
        try:
            text = content.decode('utf-8')

            return {
                'text': text,
                'structure': {
                    'type': 'text'
                }
            }

        except Exception as e:
            logger.error("Text processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process text: {str(e)}")

    async def _process_image(self, content: bytes) -> Dict[str, Any]:
        """
        Process image files with OCR
        """
        try:
            # Open image
            image = Image.open(io.BytesIO(content))

            # Perform OCR
            text = pytesseract.image_to_string(image)

            # Get image info
            image_info = {
                'format': image.format,
                'mode': image.mode,
                'size': image.size,
                'width': image.width,
                'height': image.height
            }

            return {
                'text': text,
                'image_metadata': image_info,
                'structure': {
                    'type': 'image'
                }
            }

        except Exception as e:
            logger.error("Image processing failed", error=str(e))
            raise DocumentProcessingError(f"Failed to process image: {str(e)}")


# Global instance
document_processor = DocumentProcessor()
